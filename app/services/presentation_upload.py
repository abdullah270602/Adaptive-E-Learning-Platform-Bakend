
import logging
import os
from pptx import Presentation
from app.database.connection import PostgresConnection
from app.database.slides_queries import create_slide_query
from app.services.minio_client import MinIOClientContext, save_file_to_minio
from app.services.pdf_converter import convert_to_pdf


logger = logging.getLogger(__name__)


async def process_uploaded_slides(tmp_path: str, original_filename: str, user_id: str):
    """ Processes a slides uploaded by the user """
    try:
        prs = Presentation(tmp_path)
        total_slides = len(prs.slides)
        has_notes = any(slide.has_notes_slide for slide in prs.slides) # TODO might remove this later if too slow

        pdf_path = await convert_to_pdf(tmp_path)
        
        pdf_filename = os.path.basename(pdf_path)
        s3_key = f"user_uploads/{user_id}/{pdf_filename}"

        with MinIOClientContext() as s3:
            await save_file_to_minio(s3, pdf_path, s3_key)

        with PostgresConnection() as conn:
            presentation_id = create_slide_query(
                conn, user_id, original_filename, original_filename, s3_key,
                total_slides=total_slides,
                has_speaker_notes=has_notes
            )

        os.remove(pdf_path)

        return {
            "book_metadata": None,
            "presentation_metadata": {
                "type": "presentation",
                "presentation_id": presentation_id,
                "title": original_filename,
                "slides": total_slides,
                "has_notes": has_notes,
            },
            "note_metadata": None,
        }
    
    except Exception as e:
        import traceback; traceback.print_exc();
        logger.error(f"Error processing uploaded slides: {e}")
        raise
