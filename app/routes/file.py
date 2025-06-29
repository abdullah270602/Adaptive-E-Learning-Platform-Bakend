import logging
import os
import traceback
import uuid
from pptx import Presentation
from fastapi import APIRouter, Depends, UploadFile, File, Form, HTTPException, status
from app.auth.dependencies import get_current_user
from app.database.book_queries import (
    create_book_query,
    get_book_structure_query,
    get_books_by_user,
    get_section_content_query,
)
from app.database.connection import PostgresConnection
from app.database.slides_queries import create_slide_query, get_slides_by_user
from app.services.delete_service import delete_document_and_assets
from app.services.minio_client import MinIOClientContext
from app.services.book_processor import process_toc_pages

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/file", tags=["Files"])

# TODO Make this endpoint modular divided by document type
# and add more validation for file types and sizes
@router.post(
    "/upload",
    status_code=status.HTTP_201_CREATED,
)
async def upload_file(
    file: UploadFile = File(...),
    document_type: str = Form(...),  # "book", "presentation", "notes"
    toc_pages: str = Form(None), # TODO Add validation for this field
    current_user: str = Depends(get_current_user),
):
    try:
        if document_type == "notes":
            return {
                "type": document_type,
                "notes_id": "",
                "title": "",
                "content": "Notes feature is not implemented yet.",
            }
        # Validate file extension
        ext = file.filename.split(".")[-1].lower()
        if document_type == "book" and ext != "pdf":
            raise HTTPException(status_code=400, detail="Books must be in PDF format.")
        elif document_type == "slides" and ext != "pptx":
            raise HTTPException(status_code=400, detail="Slides must be in .pptx format.")


        # Save file temporarily
        unique_name = f"{uuid.uuid4()}_{file.filename}"
        tmp_path = f"/tmp/{unique_name}"  # Use /tmp for Unix-like systems

        # tmp_path = os.path.join(os.getenv("TMP", "temp"), unique_name) # Use when on windows

        file_bytes = await file.read()  # Await first
        with open(tmp_path, "wb") as f:
            f.write(file_bytes)  # Then write synchronously

        #  Upload File to MinIO
        s3_key = f"user_uploads/{current_user}/{unique_name}"
        with MinIOClientContext() as s3:
            bucket = os.getenv("MINIO_BUCKET_NAME")
            s3.upload_file(Filename=tmp_path, Bucket=bucket, Key=s3_key)

        # Process TOC if book with TOC page range
        result = None
        book_id = unique_name.split("_")[0]

        if document_type == "book" and toc_pages:
            start_page, end_page = map(int, toc_pages.split("-"))

            with PostgresConnection() as conn:
                # Create book entry in database
                book = create_book_query(
                    conn=conn,
                    user_id=current_user,
                    book_id=book_id,
                    title=file.filename,
                    file_name=file.filename,
                    s3_key=s3_key,
                )

            result = await process_toc_pages(
                pdf_path=tmp_path,
                start_page=start_page,
                end_page=end_page,
                book_id=book_id,
                s3_key=s3_key,
            )
            result["type"] = document_type
        elif document_type == "presentation":
            
            prs = Presentation(tmp_path)
            total_slides = len(prs.slides)
            has_notes = any(slide.has_notes_slide for slide in prs.slides)

            # Save metadata in Postgres
            with PostgresConnection() as conn:
                presentation_id = create_slide_query(
                    conn=conn,
                    user_id=current_user,
                    title=file.filename,
                    original_filename=file.filename,
                    s3_key=s3_key,
                    total_slides=total_slides,
                    has_speaker_notes=has_notes,
                )

            result = {
                "type": document_type,
                "presentation_id": presentation_id,
                "title": file.filename,
                "slides": total_slides,
                "has_notes": has_notes,
            }

        else:
            logger.info(f"No TOC processing needed for document_type={document_type}")

        os.remove(tmp_path)

        return {
            "message": "Upload successful",
            "s3_key": s3_key,
            "book_metadata": result if document_type == "book" else None,
            "presentation_metadata": result if document_type == "slides" or document_type == "presentation" else None
        }
    
    
    except HTTPException as http_err:
        raise http_err

    except Exception as e:
        traceback.print_exc()
        logger.error(f"[Upload] Failed to upload: {str(e)}")
        raise HTTPException(status_code=500, detail="Upload failed.")


@router.get("/books")
async def list_user_books(
    current_user: str = Depends(get_current_user),
):
    try:
        with PostgresConnection() as conn:

            books = get_books_by_user(conn, current_user)
            return {"books": books}
    except Exception as e:
        logger.error(
            f"[List Books] Error retrieving books for user {current_user}: {str(e)}"
        )
        raise HTTPException(status_code=500, detail="Failed to retrieve user books")


@router.get("/book-structure/{book_id}", status_code=status.HTTP_200_OK)
async def get_book_structure(
    book_id: uuid.UUID, current_user: str = Depends(get_current_user)
):
    try:
        with PostgresConnection() as conn:
            return get_book_structure_query(conn, book_id)
    except Exception as e:
        traceback.print_exc();
        raise HTTPException(status_code=500, detail="Failed to retrieve book structure")


@router.get("/section/{section_id}", status_code=status.HTTP_200_OK)
async def get_section_content(
    section_id: uuid.UUID, current_user: str = Depends(get_current_user)
):
    try:
        with PostgresConnection() as conn:
            return get_section_content_query(conn, section_id)
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(
            status_code=500, detail="Failed to retrieve section content"
        )


@router.get("/slides", status_code=status.HTTP_200_OK)
async def list_user_slides(current_user: str = Depends(get_current_user)):
    try:
        presentations = []
        with PostgresConnection() as conn:
            presentations = get_slides_by_user(conn, current_user)

        return {"presentations": presentations}
    
    except Exception as e:
        logging.error(f"[Slides] Failed to fetch presentations for user {current_user}: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to retrieve presentation list")
    

@router.delete("/delete/{document_type}/{document_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_file(
    document_type: str,
    document_id: str,
    current_user: str = Depends(get_current_user),
):
    """
    Deletes any document type: book, slides, notes
    """
    try:
        deleted = delete_document_and_assets(
            document_type=document_type,
            document_id=document_id,
            user_id=current_user,
        )
        if not deleted:
            raise HTTPException(status_code=404, detail=f"{document_type.capitalize()} not found or not owned by user")

        return {"message": f"{document_type.capitalize()} deleted successfully"}

    except Exception as e:
        logging.error(f"[Delete] Failed to delete {document_type}: {e}")
        raise HTTPException(status_code=500, detail="Deletion failed.")
